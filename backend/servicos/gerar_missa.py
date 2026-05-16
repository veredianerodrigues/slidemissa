#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
gerar_missa.py - Logica de geracao de slides da Missa Catolica
"""

import re
import math
import unicodedata
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.oxml.ns import qn as _qn
from pptx.oxml.ns import qn as _pptx_qn


# -- Constantes ---------------------------------------------------------------

TEXT_COLOR = RGBColor(0x00, 0x00, 0x00)
BG_COLOR   = RGBColor(0xFF, 0xFF, 0xFF)

BOX_LEFT_CM   = 0.0
BOX_TOP_CM    = 0.0
BOX_WIDTH_CM  = 25.4
BOX_HEIGHT_CM = 19.05

FONT_NAME = "Verdana"
MAX_LINES_PER_SLIDE = 8
FIXED_FONT_SIZE = 55
HEADER_RE    = re.compile(r'^\[(.+)\]\s*$')
BRACKET_RE   = re.compile(r'\[([^\]]+)\]')

TITLE_ALIASES = {
    'gloria': 'hino de louvor',
    'glória': 'hino de louvor',
    'canto final': 'final',
}

TITULOS_PADRAO = [
    'CANTO DE ENTRADA',
    'CANTO DO ATO PENITENCIAL',
    'CANTO DO GLÓRIA',
    'CANTO DE ACLAMAÇÃO',
    'CANTO DO OFERTÓRIO',
    'SANTO',
    'CANTO DE COMUNHÃO',
    'CANTO FINAL',
]

KEYWORDS_MAP = {
    'CANTO DE ENTRADA': [
        'entrada', 'inicio', 'início', 'abertura', 'procissao', 'procissão', 'canto inicial'
    ],
    'CANTO DO ATO PENITENCIAL': [
        'penitencial', 'ato penitencial', 'canto penitencial'
    ],
    'CANTO DO GLÓRIA': [
        'gloria', 'glória', 'hino de louvor'
    ],
    'CANTO DE ACLAMAÇÃO': [
        'aclamacao', 'aclamação', 'evangelho'
    ],
    'CANTO DO OFERTÓRIO': [
        'ofertorio', 'ofertório', 'ofertas', 'apresentacao das oferendas', 'apresentação das oferendas'
    ],
    'SANTO': [
        'hosana', 'santo santo', 'santo, santo, santo'
    ],
    'CANTO DE COMUNHÃO': [
        'comunhao', 'comunhão', 'hino de comunhão'
    ],
    'CANTO FINAL': [
        'final', 'saida', 'saída', 'envio', 'despedida', 'canto final'
    ],
}


def normalize(text):
    text = text.lower()
    text = re.sub(r'\bcanto\s+(?:de|da|do|d[oa]s?)\s+', '', text)

    for alias, canonical in TITLE_ALIASES.items():
        text = re.sub(r'\b' + re.escape(alias) + r'\b', canonical, text)

    text = ''.join(
        c for c in unicodedata.normalize('NFD', text)
        if unicodedata.category(c) != 'Mn'
    )
    return re.sub(r'[\[\]\s]+', ' ', text).strip()


def _normalize_kw(text):
    text = text.lower()
    return ''.join(
        c for c in unicodedata.normalize('NFD', text)
        if unicodedata.category(c) != 'Mn'
    )


def _para_is_bold(para):
    """Detecta negrito via runs explícitos ou hierarquia de estilos do parágrafo."""
    runs = para.runs
    if not runs:
        return False
    bold_votes = 0
    for run in runs:
        if run.bold is True:
            bold_votes += 1
        elif run.bold is None:
            # Negrito não definido no run — verifica a hierarquia de estilos
            try:
                style = para.style
                while style:
                    if style.font.bold is True:
                        bold_votes += 1
                        break
                    style = style.base_style
            except Exception:
                pass
    return bold_votes > len(runs) / 2


def _match_keyword(text_norm):
    for title, keywords in KEYWORDS_MAP.items():
        for kw in keywords:
            if _normalize_kw(kw) in text_norm:
                return title
    return None


def _is_section_title(line):
    """Linha de título: toda em maiúsculas, curta (<= 60 chars), sem pontuação de frase.

    Aceita títulos com dois-pontos no final (ex: 'ENTRADA:', 'GLÓRIA:').
    """
    line = line.lstrip('* ').strip()
    if not line or len(line) > 60:
        return False
    # Dois-pontos no final é formatação de título, não conteúdo — remove antes de validar
    if line.endswith(':'):
        line = line[:-1].strip()
    if not line:
        return False
    if any(c in line for c in '.!?,;'):
        return False
    return line.upper() == line and len(line) > 4 and any(c.isalpha() for c in line)


def _extrair_secoes_docx(doc):
    """Lê parágrafos do Document e retorna lista de seções brutas com metadados.

    Títulos são detectados em qualquer linha — não exige linha em branco anterior.
    Linhas em branco dentro de uma seção são preservadas como separadores de slide.
    """
    raw_lines = []
    for para in doc.paragraphs:
        full_text = para.text
        if not full_text.strip():
            raw_lines.append('')
            continue
        is_bold = _para_is_bold(para)
        for sub in full_text.split('\n'):
            text = sub.strip()
            if text:
                raw_lines.append(f'* {text}' if is_bold else text)

    sections_raw = []
    current_raw_title = None
    current_lines = []

    for line in raw_lines:
        clean = line.lstrip('* ').strip()

        if clean and _is_section_title(clean):
            if current_raw_title is not None:
                sections_raw.append({
                    'titulo_original': current_raw_title,
                    'lines': current_lines,
                })
            # Remove dois-pontos de formatação do título (ex: "ENTRADA:" → "ENTRADA")
            current_raw_title = clean[:-1].strip() if clean.endswith(':') else clean
            current_lines = []
        else:
            if current_raw_title is not None:
                # preserva linhas em branco como separadores de slide
                if line.strip() or current_lines:
                    current_lines.append(line)

    if current_raw_title is not None:
        sections_raw.append({
            'titulo_original': current_raw_title,
            'lines': current_lines,
        })

    # Remove linhas em branco no início/fim de cada seção
    for s in sections_raw:
        lines = s['lines']
        while lines and not lines[0].strip():
            lines.pop(0)
        while lines and not lines[-1].strip():
            lines.pop()

    # Aplica mapeamento por keyword e mapeamento posicional
    for s in sections_raw:
        raw = s['titulo_original']
        matched = _match_keyword(_normalize_kw(raw))
        s['title'] = matched if matched else raw
        s['auto'] = not matched

    if 6 <= len(sections_raw) <= len(TITULOS_PADRAO):
        for i, s in enumerate(sections_raw):
            if s['auto']:
                s['title'] = TITULOS_PADRAO[i]

    return sections_raw


def validar_docx(docx_path):
    """Valida o DOCX antes de gerar o PPTX. Retorna erros, avisos e resumo das seções."""
    try:
        doc = Document(docx_path)
        sections_raw = _extrair_secoes_docx(doc)
    except Exception as e:
        return {'valido': False, 'erros': [f'Não foi possível ler o arquivo: {e}'], 'avisos': [], 'secoes': []}

    erros = []
    avisos = []
    secoes_info = []

    if not sections_raw:
        erros.append(
            'Nenhuma seção (título) encontrada. '
            'Os títulos devem estar em letras MAIÚSCULAS (ex: CANTO DE ENTRADA).'
        )
        return {'valido': False, 'erros': erros, 'avisos': avisos, 'secoes': []}

    for s in sections_raw:
        title = s['title']
        raw_title = s['titulo_original']
        lines = s['lines']
        reconhecido = not s['auto']

        tem_conteudo = any(l.strip() for l in lines)
        tem_negrito = any(l.startswith('* ') for l in lines)
        total_linhas = sum(1 for l in lines if l.strip())

        linhas = [
            {
                'texto': l[2:] if l.startswith('* ') else l,
                'negrito': l.startswith('* '),
            }
            for l in lines
        ]

        secoes_info.append({
            'titulo': title,
            'titulo_original': raw_title,
            'titulo_reconhecido': reconhecido,
            'tem_conteudo': tem_conteudo,
            'tem_negrito': tem_negrito,
            'total_linhas': total_linhas,
            'linhas': linhas,
        })

        if not tem_conteudo:
            erros.append(f'Seção "{raw_title}" está vazia — nenhum conteúdo encontrado após o título.')
        else:
            if not tem_negrito:
                avisos.append(
                    f'Seção "{title}" não tem texto em negrito. '
                    'Refrões devem estar em negrito (selecione e pressione Ctrl+B).'
                )
            if not reconhecido:
                avisos.append(
                    f'Título "{raw_title}" não é um título litúrgico padrão — '
                    'verifique se está escrito corretamente.'
                )

    return {
        'valido': len(erros) == 0,
        'erros': erros,
        'avisos': avisos,
        'secoes': secoes_info,
    }


def parse_docx(docx_path, log=None):
    """Lê um DOCX e retorna seções no mesmo formato que parse_txt() retornava."""
    doc = Document(docx_path)
    sections_raw = _extrair_secoes_docx(doc)  # já aplica mapeamento posicional

    # Converte para o formato interno de blocos usado pelo gerador
    sections = {}
    for s in sections_raw:
        title = s['title']
        key = normalize(title)

        # Agrupa por linha em branco (separador de slide explícito)
        grupos = []
        grupo_atual = []
        for line in s['lines']:
            stripped = line.strip()
            if not stripped:
                if grupo_atual:
                    grupos.append(grupo_atual)
                    grupo_atual = []
            else:
                is_bold = stripped.startswith('* ')
                text = stripped[2:] if is_bold else stripped
                if text:
                    grupo_atual.append({'text': text, 'bold': is_bold})
        if grupo_atual:
            grupos.append(grupo_atual)

        expanded = []
        for grupo in grupos:
            for sub in split_block(grupo):
                expanded.extend(auto_split_large(sub, 48))

        if expanded:
            sections[key] = {'titulo': title, 'blocos': expanded}

    if log:
        log('Secoes encontradas: ' + str(len(sections)))
        for key, data in sections.items():
            log('  [' + data['titulo'] + ']: ' + str(len(data['blocos'])) + ' slide(s)')

    return sections


def split_block(block_lines):
    """Agrupa todas as linhas em um único bloco, mantendo negrito por linha."""
    if not block_lines:
        return []
    return [{'lines': [{'text': l['text'], 'bold': l['bold']} for l in block_lines]}]


def estimate_lines_needed(text_lines, font_size_pt, bold=False, box_width_cm=BOX_WIDTH_CM):
    box_width_emu  = box_width_cm * 360000
    char_width_factor = 0.65 if bold else 0.60
    avg_char_w_emu = font_size_pt * char_width_factor * 12700
    chars_per_line = box_width_emu / avg_char_w_emu
    return sum(max(1, math.ceil(len(ln) / chars_per_line)) for ln in text_lines)


def estimate_max_lines(font_size_pt, box_height_cm=BOX_HEIGHT_CM):
    box_height_emu  = box_height_cm * 360000
    line_height_emu = font_size_pt * 1.25 * 12700
    return int(box_height_emu / line_height_emu)


def auto_split_large(block, font_min, box_width_cm=BOX_WIDTH_CM, box_height_cm=BOX_HEIGHT_CM):
    lines = block['lines']  # list of {'text': ..., 'bold': ...}
    n = len(lines)

    if n == 0:
        return []

    font_est = FIXED_FONT_SIZE if FIXED_FONT_SIZE else font_min
    text_lines = [l['text'] for l in lines]
    has_bold = any(l['bold'] for l in lines)
    max_lines = estimate_max_lines(font_est, box_height_cm=box_height_cm)

    if n == 1:
        if estimate_lines_needed(text_lines, font_est, bold=has_bold, box_width_cm=box_width_cm) <= max_lines:
            return [block]
        words = lines[0]['text'].split()
        if len(words) <= 1:
            return [block]
        mid = math.ceil(len(words) / 2)
        left  = {'lines': [{'text': ' '.join(words[:mid]), 'bold': lines[0]['bold']}]}
        right = {'lines': [{'text': ' '.join(words[mid:]), 'bold': lines[0]['bold']}]}
        return (auto_split_large(left, font_min, box_width_cm, box_height_cm) +
                auto_split_large(right, font_min, box_width_cm, box_height_cm))

    if n <= MAX_LINES_PER_SLIDE and estimate_lines_needed(text_lines, font_est, bold=has_bold, box_width_cm=box_width_cm) <= max_lines:
        return [block]

    mid = math.ceil(n / 2)
    left  = {'lines': lines[:mid]}
    right = {'lines': lines[mid:]}
    return auto_split_large(left, font_min, box_width_cm, box_height_cm) + auto_split_large(right, font_min, box_width_cm, box_height_cm)


def slide_full_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
    return ' '.join(parts)


def slide_normalized_key(slide):
    txt = slide_full_text(slide).strip()
    if not txt or len(txt) > 80:
        return None
    clean = re.sub(r'\s*\|\s*', ' ', txt).strip()
    return normalize(clean)


def get_bg_color(slide):
    try:
        fill = slide.background.fill
        if fill.type is not None:
            return fill.fore_color.rgb
    except Exception:
        pass
    return BG_COLOR


def detect_slide_format(prs):
    slide_width_cm  = prs.slide_width  / 360000
    slide_height_cm = prs.slide_height / 360000
    return slide_width_cm, slide_height_cm


def choose_font_size(text_lines, font_min=48, font_max=60, box_width_cm=BOX_WIDTH_CM, box_height_cm=BOX_HEIGHT_CM):
    for size in range(font_max, font_min - 2, -2):
        if estimate_lines_needed(text_lines, size, box_width_cm=box_width_cm) <= estimate_max_lines(size, box_height_cm=box_height_cm):
            return size
    return font_min


def create_song_slide(prs, block, bg_color, box_width_cm, box_height_cm, font_min=48, font_max=60, fixed_font_size=None):
    lines      = block['lines']  # list of {'text': ..., 'bold': ...}
    text_lines = [l['text'] for l in lines]
    font_size  = fixed_font_size if fixed_font_size else choose_font_size(text_lines, font_min, font_max, box_width_cm, box_height_cm)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for shape in list(slide.shapes):
        if shape.is_placeholder or (hasattr(shape, 'text') and not shape.text.strip()):
            sp = shape.element
            sp.getparent().remove(sp)

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    txBox = slide.shapes.add_textbox(
        Cm(BOX_LEFT_CM), Cm(BOX_TOP_CM),
        Cm(box_width_cm), Cm(box_height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    body_pr = tf._txBody.find(_pptx_qn('a:bodyPr'))
    if body_pr is not None:
        body_pr.set('anchor', 'ctr')

    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text           = line['text']
        run.font.name      = FONT_NAME
        run.font.size      = Pt(font_size)
        run.font.bold      = line['bold']
        run.font.color.rgb = TEXT_COLOR

    return slide


def move_slide(prs, from_idx, to_idx):
    if from_idx == to_idx:
        return
    sldIdLst = prs.slides._sldIdLst
    all_ids  = list(sldIdLst)
    sldId    = all_ids.pop(from_idx)
    all_ids.insert(to_idx, sldId)
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in all_ids:
        sldIdLst.append(el)


def _secoes_to_dict(secoes_editadas):
    """Converte o formato JSON da API (linhas com negrito) para o dict interno de seções.

    Linhas com texto vazio funcionam como separadores explícitos de slide.
    """
    sections = {}
    for s in secoes_editadas:
        titulo = s['titulo']
        key = normalize(titulo)

        grupos = []
        grupo_atual = []
        for l in s['linhas']:
            if not l['texto'].strip():
                if grupo_atual:
                    grupos.append(grupo_atual)
                    grupo_atual = []
            else:
                grupo_atual.append({'text': l['texto'], 'bold': l['negrito']})
        if grupo_atual:
            grupos.append(grupo_atual)

        expanded = []
        for grupo in grupos:
            for sub in split_block(grupo):
                expanded.extend(auto_split_large(sub, 48))

        if expanded:
            sections[key] = {'titulo': titulo, 'blocos': expanded}
    return sections


def _inserir_cantos(prs, cantos, output_path, fonte_min=48, fonte_max=60, log=None):
    """Busca placeholders no PPTX, insere slides de cantos e salva."""
    if log is None:
        log = lambda msg: None

    log('Total de slides: ' + str(len(prs.slides)))

    placeholders = []
    for i, slide in enumerate(prs.slides):
        key = slide_normalized_key(slide)
        if key and key in cantos:
            placeholders.append((i, key))
            log('  Slide ' + str(i + 1).zfill(2) + ': [' + cantos[key]['titulo'] + ']')

    if not placeholders:
        return False, 'Nenhum titulo do DOCX foi encontrado no PPTX.'

    box_width, box_height = detect_slide_format(prs)
    log(f'Dimensões do slide: {box_width:.2f} x {box_height:.2f} cm')

    total_inserted = 0
    for ph_idx, sec_key in reversed(placeholders):
        data     = cantos[sec_key]
        blocos   = data['blocos']
        bg_color = get_bg_color(prs.slides[ph_idx])
        log('Inserindo ' + str(len(blocos)) + ' slide(s) para [' + data['titulo'] + ']')
        slide_idx = 0
        for block in blocos:
            sub_blocos = auto_split_large(block, fonte_min, box_width, box_height)
            for sub_block in sub_blocos:
                create_song_slide(prs, sub_block, bg_color, box_width, box_height, fonte_min, fonte_max, fixed_font_size=FIXED_FONT_SIZE)
                move_slide(prs, len(prs.slides) - 1, ph_idx + 1 + slide_idx)
                slide_idx += 1
        total_inserted += slide_idx

    prs.save(output_path)
    msg = str(total_inserted) + ' slides inseridos. Total final: ' + str(len(prs.slides)) + ' slides.'
    log(msg)
    return True, msg


def gerar(docx_path, pptx_path, output_path, fonte_min=48, fonte_max=60, log=None):
    if log is None:
        log = lambda msg: None

    try:
        log('Lendo cantos: ' + docx_path)
        cantos = parse_docx(docx_path, log)

        log('Carregando ritual: ' + pptx_path)
        prs = Presentation(pptx_path)

        return _inserir_cantos(prs, cantos, output_path, fonte_min, fonte_max, log)

    except Exception as e:
        return False, 'Erro: ' + str(e)


def gerar_de_secoes(secoes_editadas, pptx_path, output_path, fonte_min=48, fonte_max=60, log=None):
    """Gera o PPTX a partir de seções já editadas (sem releitura do DOCX)."""
    if log is None:
        log = lambda msg: None

    try:
        cantos = _secoes_to_dict(secoes_editadas)
        log('Seções recebidas: ' + str(len(cantos)))

        log('Carregando ritual: ' + pptx_path)
        prs = Presentation(pptx_path)

        return _inserir_cantos(prs, cantos, output_path, fonte_min, fonte_max, log)

    except Exception as e:
        return False, 'Erro: ' + str(e)
